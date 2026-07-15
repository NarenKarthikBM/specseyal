#!/usr/bin/env python3
"""render.py -- markdown -> pptx, the deterministic transform (T015).

This is the integrating script for the deck-render extension: it maps a
defense deck's ordered `Block` model (`deck_md.parse()`) onto pptx slides
per the T1-T10 rules (`specs/006-deck-render/data-model.md` Sec 4), resolves
which deck(s) to render (`profile_key.resolve_deck_render()`, FR-016), and
disclose a per-deck outcome with the right exit code
(`specs/006-deck-render/contracts/commands.md`).

Contract surface (normative -- see the docs, not re-derived here):
  - CLI: `render.py [technical|overview|both] [--feature <dir>]
    [--validate-profile]` -- contracts/commands.md Sec 1.
  - Behavior + invariants I1-I7: contracts/commands.md Sec 2.
  - Exit codes 0/2/3/4 + the `both` outcome matrix: contracts/commands.md
    Sec 4.
  - The derived-render stamp: data-model.md Sec 3.
  - T1-T10 slide mapping + the FR-002 fidelity allowlist: data-model.md
    Sec 4.
  - Render outcome shape (O1-O5): data-model.md Sec 5.
  - Phase B invariants I-B1..I-B6: plan.md.

Design decisions load-bearing enough to restate here:

  - **`python-pptx` is imported lazily, INSIDE the render path, never at
    module top level** (FR-015, R2). Installing this extension must never
    require a presentation toolchain, and an absent toolchain must degrade
    to a disclosed per-deck failure, never a crash. Every name this module
    needs from `pptx` is imported inside `_render_deck()`'s try block --
    nothing pptx-shaped appears above that point.
  - **The stamp's SHA is sha256 of the source markdown's BYTES, not git's
    commit SHA** (R3): a commit SHA fails closed on the deck's routinely
    uncommitted state at gate time and cannot detect a working-tree edit
    (which would silently break SC-007). The full 64-hex digest is
    embedded on every slide, never truncated (I-B5) -- its length is what
    makes it self-evidently not a 40-hex git SHA.
  - **The write is atomic** (I-B3/O5): the pptx zip lands in a temp file
    inside `renders/`, and only `os.replace()`s over the final path on
    full success. No target is ever pre-deleted, so a failed render never
    disturbs a prior good one.
  - **Per-deck isolation** (I-B2): the loop over selected decks catches
    every failure locally: an exception rendering one deck never prevents
    the next deck from being attempted and reported.
  - **FRESH/STALE is a stateless read-and-compare** (I-B6): when a prior
    render already sits at the target path, this module reads ITS embedded
    stamp back out via plain `zipfile` + `xml.etree` (no `pptx` needed to
    READ, only to write) and compares its SHA to the current source's. No
    state file is written anywhere (principle 3).
  - **The renderer never authors** (T10). Anything the deterministic
    mapping cannot place -- more than one H1, content before the H1, a
    markdown construct `deck_md.py` itself rejects -- is a `failed`
    outcome, never a silent simplification.

Not this module's job: parsing markdown (`deck_md.py`), or resolving/
validating `deck_render` out of `profile.yaml` (`profile_key.py`). Both are
imported, never re-implemented.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from xml.etree import ElementTree

# Make the sibling deck_md.py / profile_key.py importable regardless of how
# this script is loaded -- `python3 render.py ...` auto-prepends this
# file's directory to sys.path, but a test harness importing this module
# in-process from a different cwd does not get that for free. Mirrors
# extensions/workforce/extension/scripts/validate-categorization.py's
# identical guard.
sys.path.insert(0, str(Path(__file__).resolve().parent))

import deck_md  # noqa: E402
from profile_key import DECK_RENDER_ENUM, ProfileKeyError, resolve_deck_render  # noqa: E402

# ---------------------------------------------------------------------------
# Constants -- exit codes, the derived-render stamp text, layout budgets.
# None of this is pptx-shaped; it is safe at module scope.
# ---------------------------------------------------------------------------

#: contracts/commands.md Sec 4.
EXIT_OK = 0
EXIT_PARTIAL = 2
EXIT_INVALID = 3
EXIT_ALL_FAILED = 4

#: The two decks this extension ever writes a file for, in the fixed,
#: deterministic order every disclosure and every `both` expansion uses.
#: Derived from `profile_key.DECK_RENDER_ENUM` (the SSOT) rather than
#: re-typed, per the task's "do NOT re-derive the enum" instruction.
RENDERABLE_DECKS = tuple(v for v in DECK_RENDER_ENUM if v not in ("none", "both"))

#: The CLI's explicit-selection choices (commands.md Sec 1): every enum
#: member except `none` -- there is no `render.py none` invocation; "render
#: nothing" is what an absent argument + an absent/`none` profile already
#: means (commands.md Sec 2 step 3).
_CLI_DECK_CHOICES = tuple(v for v in DECK_RENDER_ENUM if v != "none")

#: Per-deck outcomes (data-model.md Sec 5).
OUTCOME_RENDERED = "rendered"
OUTCOME_FAILED = "failed"
OUTCOME_SKIPPED = "skipped"

# --- the derived-render stamp (data-model.md Sec 3) -------------------------

STAMP_DECLARATION = "Derived render — NOT the artifact of record."
STAMP_POINTER = (
    "The markdown at the path above is what the council reviewed and what the gate binds."
)

# --- T7 overflow: a fixed, deterministic per-slide line budget --------------
#
# A module constant (not a computed/measured value) is what makes the
# `(cont.)` split deterministic (T7): same input bytes -> same slide count,
# always, on any host, with or without the real pptx toolchain to measure
# actual glyph widths against. The estimate below is a content-only
# heuristic (character count / wrap width, or literal row/line counts for
# tables and code) -- not real typography -- precisely so it never depends
# on a font metric that could differ across a viewer or a host.
LINE_BUDGET_PER_SLIDE = 20
_WRAP_CHARS_PER_LINE = 88

# --- slide layout constants (plain floats, inches; pptx.util.Inches() is
# applied only inside the pptx-dependent helpers below) ---------------------
_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5
_BODY_TOP_IN = 1.3
_BODY_LEFT_IN = 0.6
_BODY_WIDTH_IN = 12.1
_ITEM_GAP_IN = 0.12
_LINE_HEIGHT_IN = 0.28
_CODE_LINE_HEIGHT_IN = 0.2
_TABLE_ROW_HEIGHT_IN = 0.35
_FOOTER_TOP_IN = 6.95
_FOOTER_LEFT_IN = 0.4
_FOOTER_WIDTH_IN = 11.6


# ---------------------------------------------------------------------------
# Errors
# ---------------------------------------------------------------------------


class RenderInputError(Exception):
    """An invalid-input condition at the CLI boundary: an unresolvable
    feature directory. Maps to exit 3 (contracts/commands.md Sec 4), the
    same family as `profile_key.ProfileKeyError`.
    """


class RenderError(Exception):
    """T10: the renderer never authors. Raised when the deterministic
    mapping has no defined home for something it was handed -- more than
    one H1, content before the H1 -- so the caller can turn it into a
    disclosed per-deck `failed` outcome (degrade + disclose) rather than
    silently reinterpreting the deck.
    """


# ---------------------------------------------------------------------------
# Result / logical-slide data shapes
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class Result:
    """One deck's render outcome (data-model.md Sec 5)."""

    deck: str
    outcome: str  # OUTCOME_RENDERED | OUTCOME_FAILED | OUTCOME_SKIPPED
    path: Path | None = None
    reason: str | None = None


@dataclass
class LogicalSlide:
    """One un-split slide's worth of content, before T7 overflow is
    applied: either the `Preamble` slide or one H2 section. `items` is a
    list of `deck_md.Block`s in source order (T5); HR blocks and the H1
    itself are never present here -- they never had a body slot to begin
    with (T1/T8).
    """

    kind: str  # "preamble" | "section"
    title: str
    items: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# Feature-directory + selection resolution (commands.md Sec 2 step 1-2)
# ---------------------------------------------------------------------------


def _resolve_feature_dir(explicit: str | None) -> Path:
    """`--feature <dir>`, else `.specify/feature.json`'s `feature_directory`
    (used directly as the feature dir -- it is already a path like
    `specs/006-deck-render`, not a spec-id to reconstruct one from), else
    the current git branch (`specs/<branch>`). Raises `RenderInputError`
    (exit 3) only when none of the three resolves -- an unresolvable
    feature directory, per commands.md Sec 4.
    """
    if explicit:
        return Path(explicit)

    feature_json = Path(".specify") / "feature.json"
    if feature_json.is_file():
        try:
            data = json.loads(feature_json.read_text(encoding="utf-8"))
        except (OSError, ValueError):
            data = None
        if isinstance(data, dict):
            raw = data.get("feature_directory")
            if isinstance(raw, str) and raw.strip():
                return Path(raw.strip())

    try:
        proc = subprocess.run(
            ["git", "rev-parse", "--abbrev-ref", "HEAD"],
            capture_output=True,
            text=True,
            timeout=10,
        )
    except (OSError, subprocess.SubprocessError) as exc:
        raise RenderInputError(
            "cannot resolve a feature directory: no --feature given, no usable "
            f".specify/feature.json, and git is unavailable ({exc})"
        ) from exc

    branch = proc.stdout.strip() if proc.returncode == 0 else ""
    if not branch or branch == "HEAD":
        raise RenderInputError(
            "cannot resolve a feature directory: no --feature given, no usable "
            ".specify/feature.json, and the current git ref is not a branch "
            "(detached HEAD, or not a git repo)"
        )
    return Path("specs") / branch


def _expand_selection(value: str) -> tuple[str, ...]:
    """`none` -> nothing selected; `both` -> both renderable decks, in the
    fixed canonical order; a single deck name -> just that one.
    """
    if value == "none":
        return ()
    if value == "both":
        return RENDERABLE_DECKS
    return (value,)


# ---------------------------------------------------------------------------
# The derived-render stamp -- computing it, and reading one back (I-B6)
# ---------------------------------------------------------------------------

#: Matches the full, never-truncated 64-hex sha256 embedded on every slide
#: (I-B5). Used only to read a PRIOR render's own stamp back out, via plain
#: zipfile + XML -- no `pptx` import needed to read, only to write.
_SHA256_RE = re.compile(r"\b([0-9a-f]{64})\b")

_DRAWINGML_T_TAG = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"


def _read_embedded_sha256(pptx_path: Path) -> str | None:
    """Read a prior render's own embedded source-SHA back out of its title
    slide (`ppt/slides/slide1.xml`), stdlib-only (`zipfile` +
    `xml.etree.ElementTree`) -- deliberately independent of `python-pptx`,
    per I-B6 ("stdlib zip+XML is fine and needs no python-pptx to READ").
    Returns `None` for anything that isn't a readable render of ours (a
    missing/foreign/corrupt file at that path) -- FRESH/STALE is simply not
    printed in that case, rather than risking a misleading verdict.
    """
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            try:
                slide_xml = archive.read("ppt/slides/slide1.xml")
            except KeyError:
                return None
    except (OSError, zipfile.BadZipFile):
        return None

    try:
        root = ElementTree.fromstring(slide_xml)
    except ElementTree.ParseError:
        return None

    combined = "\n".join(t.text or "" for t in root.iter(_DRAWINGML_T_TAG))
    match = _SHA256_RE.search(combined)
    return match.group(1) if match else None


def _fresh_stale_verdict(target_path: Path, current_sha: str) -> str | None:
    """A stateless read-and-compare (I-B6, S13): `None` when there is
    nothing to compare against (no prior render at the target path, or an
    unreadable one); otherwise `"FRESH"` or `"STALE"` against the current
    source's sha256.
    """
    if not target_path.exists():
        return None
    existing_sha = _read_embedded_sha256(target_path)
    if existing_sha is None:
        return None
    return "FRESH" if existing_sha == current_sha else "STALE"


# ---------------------------------------------------------------------------
# T9 -- inline markdown stripped to plain text / styled runs
# ---------------------------------------------------------------------------
#
# `**bold**` and `` `code` `` are the two markers data-model.md Sec 4 names
# explicitly. A bare single-`*italic*` span (used for real in the corpus --
# e.g. the scope-note line every deck opens with) is not named by T9, but
# "the text is preserved exactly, only the markdown syntax is consumed" is
# stated as the general rule, not scoped to just the two named markers --
# leaving its literal asterisks in the render would be exactly the kind of
# stray markdown syntax T9 exists to consume. All three are stripped here;
# only bold and code additionally carry distinct run-level formatting
# (italic gets a real italic run, which is formatting, not new text).

_INLINE_RE = re.compile(
    r"\*\*(?P<bold>[^*]+?)\*\*"
    r"|`(?P<code>[^`]+?)`"
    r"|\*(?P<italic>[^*]+?)\*"
)


def _tokenize_inline(text: str) -> list[tuple[str, str]]:
    """Split `text` into `(segment_text, style)` pairs, `style` one of
    `"plain" | "bold" | "code" | "italic"`. The concatenation of every
    `segment_text` equals `text` with only the markdown marker characters
    removed -- no content is dropped or added (T9/T10).
    """
    segments: list[tuple[str, str]] = []
    pos = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > pos:
            segments.append((text[pos : m.start()], "plain"))
        if m.group("bold") is not None:
            segments.append((m.group("bold"), "bold"))
        elif m.group("code") is not None:
            segments.append((m.group("code"), "code"))
        else:
            segments.append((m.group("italic"), "italic"))
        pos = m.end()
    if pos < len(text):
        segments.append((text[pos:], "plain"))
    if not segments:
        segments.append(("", "plain"))
    return segments


def _plain_text(text: str) -> str:
    """The fully-stripped plain text of `text` -- used for H1/H2 headings,
    where T2 says the slide title carries the H2's text with "inline
    emphasis stripped to plain text" (no run-level bold/italic/monospace
    kept, unlike body content).
    """
    return "".join(seg for seg, _style in _tokenize_inline(text))


# ---------------------------------------------------------------------------
# T1-T10 -- block model -> logical slides -> T7-split physical slides
# ---------------------------------------------------------------------------


def _to_logical_slides(blocks: list) -> list[LogicalSlide]:
    """T1-T5, T8: the deck's `Block`s, in source order, folded into
    `Preamble` + one-per-H2 logical slides (not yet T7-split). Raises
    `RenderError` (T10) when the deck's H1 shape does not match what T1
    defines a home for: not exactly one H1, or non-HR content preceding it.
    """
    h1_blocks = [b for b in blocks if b.kind == deck_md.H1]
    if len(h1_blocks) != 1:
        raise RenderError(
            f"expected exactly one H1 title heading (T1), found {len(h1_blocks)}"
        )
    h1 = h1_blocks[0]

    for b in blocks:
        if b is h1:
            break
        if b.kind != deck_md.HR:
            raise RenderError(
                f"line {b.line}: content precedes the deck's H1 title heading; "
                "T1 makes the H1 slide 1 and defines no home for anything before it"
            )

    logical: list[LogicalSlide] = []
    preamble_items: list = []
    current: LogicalSlide | None = None
    seen_h2 = False

    for b in blocks:
        if b is h1 or b.kind == deck_md.HR:  # T1 handles H1 separately; T8 drops HRs
            continue
        if b.kind == deck_md.H2:
            if not seen_h2:
                if preamble_items:
                    logical.append(LogicalSlide("preamble", "Preamble", preamble_items))
                seen_h2 = True
            elif current is not None:
                logical.append(current)
            current = LogicalSlide("section", _plain_text(b.text), [])
            continue
        if not seen_h2:
            preamble_items.append(b)
        else:
            current.items.append(b)  # type: ignore[union-attr]

    if not seen_h2:
        if preamble_items:
            logical.append(LogicalSlide("preamble", "Preamble", preamble_items))
    elif current is not None:
        logical.append(current)

    return logical


def _estimate_lines(block) -> int:
    """A pure, deterministic line-count estimate for T7's overflow budget --
    content-only (character count / wrap width, or literal row/line
    counts), never a measured font metric, so the `(cont.)` split is
    reproducible on any host with or without the real toolchain.
    """
    if block.kind == deck_md.TABLE:
        return 1 + len(block.rows)
    if block.kind == deck_md.CODE:
        return max(1, len(block.lines))
    text = f"{block.number}. {block.text}" if block.kind == deck_md.NUMBERED else block.text
    return max(1, -(-len(text) // _WRAP_CHARS_PER_LINE))  # ceiling division


def _split_for_overflow(logical_slides: list[LogicalSlide]) -> list[tuple[str, list]]:
    """T7: split each logical slide into one or more `(title, items)`
    physical slides against the fixed `LINE_BUDGET_PER_SLIDE`. A single
    item whose own estimate exceeds the budget still gets its own slide
    (never split mid-item -- there is no rule for splitting a table row or
    a code block, and inventing one would be T10's silent-simplification
    failure). Every physical slide after a logical slide's first is titled
    `"<title> (cont.)"` -- the one piece of invented text T7 allowlists.
    """
    physical: list[tuple[str, list]] = []
    for ls in logical_slides:
        if not ls.items:
            physical.append((ls.title, []))
            continue

        chunk: list = []
        chunk_lines = 0
        chunk_index = 0

        for item in ls.items:
            item_lines = _estimate_lines(item)
            if chunk and chunk_lines + item_lines > LINE_BUDGET_PER_SLIDE:
                title = ls.title if chunk_index == 0 else f"{ls.title} (cont.)"
                physical.append((title, chunk))
                chunk_index += 1
                chunk, chunk_lines = [], 0
            chunk.append(item)
            chunk_lines += item_lines

        title = ls.title if chunk_index == 0 else f"{ls.title} (cont.)"
        physical.append((title, chunk))

    return physical


# ---------------------------------------------------------------------------
# pptx-dependent slide construction. Every function below is only ever
# reached AFTER `_render_deck()`'s lazy `import pptx` (and its sibling
# `from pptx... import ...` statements) has already succeeded -- nothing
# here executes, and nothing here is even parsed as reachable code, on a
# host where `python-pptx` is absent, other than these function bodies
# sitting inert until called (FR-015/R2).
# ---------------------------------------------------------------------------


@dataclass(frozen=True)
class _PptxKit:
    """The handful of `pptx` names the slide builders need, gathered once
    at the single lazy-import site and threaded through by value -- so no
    helper below needs its own `import pptx.*` statement.
    """

    Inches: object
    Pt: object
    PP_ALIGN: object
    MSO_AUTO_SIZE: object


def _add_runs(paragraph, kit: _PptxKit, segments, *, size, base_bold=False, base_italic=False) -> None:
    """Add one run per `(text, style)` segment (`_tokenize_inline`'s
    output) to `paragraph`, applying bold/italic/monospace per T9 on top
    of whatever base styling the caller wants applied to every run (e.g.
    an H3 lead line is bold throughout, on top of any inline `**bold**`
    inside it).
    """
    if not segments:
        segments = [("", "plain")]
    for seg_text, style in segments:
        run = paragraph.add_run()
        run.text = seg_text
        run.font.size = size
        run.font.bold = base_bold or (style == "bold")
        run.font.italic = base_italic or (style == "italic")
        if style == "code":
            run.font.name = "Courier New"


def _add_title_box(slide, kit: _PptxKit, text, *, top_in, height_in, font_size, align=None):
    Inches, Pt = kit.Inches, kit.Pt
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top_in), Inches(_SLIDE_WIDTH_IN - 1.0), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    return box


def _add_full_stamp(slide, kit: _PptxKit, source_display, sha_hex, *, top_in):
    """The title slide's full derived-render stamp (data-model.md Sec 3):
    declaration, source path, the full 64-hex source SHA, and the pointer
    sentence -- four elements.
    """
    Inches, Pt = kit.Inches, kit.Pt
    lines = [
        STAMP_DECLARATION,
        f"Source: {source_display}",
        f"Source SHA-256: {sha_hex}",
        STAMP_POINTER,
    ]
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top_in), Inches(_SLIDE_WIDTH_IN - 1.0), Inches(2.2)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.italic = i in (0, 3)  # declaration + pointer -- presentational only
    return box


def _add_footer_stamp(slide, kit: _PptxKit, sha_hex):
    """Every slide's abbreviated stamp (I-B5): declaration + the full,
    never-truncated 64-hex SHA -- two elements, fewer than the title
    stamp's four, but the SHA itself is exactly as long either place.
    """
    Inches, Pt = kit.Inches, kit.Pt
    box = slide.shapes.add_textbox(
        Inches(_FOOTER_LEFT_IN), Inches(_FOOTER_TOP_IN), Inches(_FOOTER_WIDTH_IN), Inches(0.45)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{STAMP_DECLARATION}  Source SHA-256: {sha_hex}"
    run.font.size = Pt(8)
    return box


def _add_slide_number(slide, kit: _PptxKit, number):
    """Structural chrome, on the FR-002 allowlist alongside the stamp and
    `(cont.)`.
    """
    Inches, Pt = kit.Inches, kit.Pt
    box = slide.shapes.add_textbox(
        Inches(_SLIDE_WIDTH_IN - 0.9), Inches(_FOOTER_TOP_IN), Inches(0.6), Inches(0.4)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(10)
    return box


def _add_table(slide, kit: _PptxKit, block, left_in, top_in, width_in) -> float:
    Inches, Pt = kit.Inches, kit.Pt
    n_cols = len(block.header)
    n_rows = 1 + len(block.rows)
    height_in = _TABLE_ROW_HEIGHT_IN * n_rows
    graphic_frame = slide.shapes.add_table(
        n_rows, n_cols, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    table = graphic_frame.table

    def _set_cell(cell, text, bold):
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_runs(p, kit, _tokenize_inline(text), size=Pt(12), base_bold=bold)

    for c, cell_text in enumerate(block.header):
        _set_cell(table.cell(0, c), cell_text, bold=True)
    for r, row in enumerate(block.rows, start=1):
        for c, cell_text in enumerate(row):
            _set_cell(table.cell(r, c), cell_text, bold=False)

    return height_in


def _add_code_box(slide, kit: _PptxKit, block, left_in, top_in, width_in) -> float:
    """T6: a monospace box, `word_wrap = False` (no reflow), one paragraph
    per source line, byte-exact -- so box-drawing and directory trees
    survive visually intact.
    """
    Inches, Pt, MSO_AUTO_SIZE = kit.Inches, kit.Pt, kit.MSO_AUTO_SIZE
    lines = block.lines if block.lines else ("",)
    height_in = max(0.3, len(lines) * _CODE_LINE_HEIGHT_IN)
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(12)
    return height_in


def _add_text_block(slide, kit: _PptxKit, block, left_in, top_in, width_in) -> float:
    """H3 (bold lead line, T4), paragraph, bullet, numbered, blockquote --
    everything that is not a table or a fenced code block.
    """
    Inches, Pt = kit.Inches, kit.Pt
    est_lines = _estimate_lines(block)
    height_in = max(0.35, est_lines * _LINE_HEIGHT_IN)
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    if block.kind == deck_md.H3:
        _add_runs(p, kit, _tokenize_inline(block.text), size=Pt(18), base_bold=True)
    elif block.kind == deck_md.BULLET:
        p.level = 1
        _add_runs(p, kit, _tokenize_inline(block.text), size=Pt(16))
    elif block.kind == deck_md.NUMBERED:
        p.level = 1
        marker_seg = [(f"{block.number}. ", "plain")]
        _add_runs(p, kit, marker_seg + _tokenize_inline(block.text), size=Pt(16))
    elif block.kind == deck_md.BLOCKQUOTE:
        p.level = 1
        _add_runs(p, kit, _tokenize_inline(block.text), size=Pt(15), base_italic=True)
    else:  # PARAGRAPH
        _add_runs(p, kit, _tokenize_inline(block.text), size=Pt(16))

    return height_in


def _add_body(slide, kit: _PptxKit, items) -> None:
    cursor_in = _BODY_TOP_IN
    for block in items:
        if block.kind == deck_md.TABLE:
            consumed = _add_table(slide, kit, block, _BODY_LEFT_IN, cursor_in, _BODY_WIDTH_IN)
        elif block.kind == deck_md.CODE:
            consumed = _add_code_box(slide, kit, block, _BODY_LEFT_IN, cursor_in, _BODY_WIDTH_IN)
        else:
            consumed = _add_text_block(slide, kit, block, _BODY_LEFT_IN, cursor_in, _BODY_WIDTH_IN)
        cursor_in += consumed + _ITEM_GAP_IN


def _assemble_presentation(
    Presentation, kit: _PptxKit, h1_title, physical_slides, source_display, sha_hex
):
    """T1-T10, assembled: title slide, then one physical slide per
    `_split_for_overflow()` entry -- every slide stamped in its footer
    (I-B5) and numbered (the allowlist).
    """
    Inches = kit.Inches
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # the built-in template's "Blank" layout

    # T1 -- the title slide: H1 + the full stamp.
    slide = prs.slides.add_slide(blank_layout)
    _add_title_box(
        slide, kit, h1_title, top_in=0.6, height_in=1.2, font_size=40, align=kit.PP_ALIGN.CENTER
    )
    _add_full_stamp(slide, kit, source_display, sha_hex, top_in=2.2)
    _add_footer_stamp(slide, kit, sha_hex)
    _add_slide_number(slide, kit, 1)

    # T2/T3/T4/T5/T6/T7 -- Preamble + each H2 section, already T7-split.
    for idx, (title, items) in enumerate(physical_slides, start=2):
        slide = prs.slides.add_slide(blank_layout)
        _add_title_box(slide, kit, title, top_in=0.3, height_in=0.8, font_size=28)
        _add_body(slide, kit, items)
        _add_footer_stamp(slide, kit, sha_hex)
        _add_slide_number(slide, kit, idx)

    return prs


# ---------------------------------------------------------------------------
# Per-deck render (commands.md Sec 2 step 4; I-B2 isolation; I-B3 atomicity)
# ---------------------------------------------------------------------------


def _render_deck(deck: str, feature_dir: Path) -> Result:
    """Render exactly one deck, never raising: every failure mode this
    function can hit becomes a `Result(outcome=OUTCOME_FAILED, ...)`
    instead, so the caller's per-deck loop is trivially isolated (I-B2) --
    an exception rendering `deck` can never prevent the next deck in the
    selection from being attempted.
    """
    source_path = feature_dir / "council" / "defense-deck" / f"{deck}.md"
    target_path = feature_dir / "renders" / f"{deck}.pptx"

    if not source_path.is_file():
        # O4: no deck present yet -- the council phase simply has not run.
        return Result(deck, OUTCOME_SKIPPED, None, "no deck present (council phase not run yet)")

    try:
        source_bytes = source_path.read_bytes()
    except OSError as exc:
        return Result(deck, OUTCOME_FAILED, None, f"cannot read source markdown: {exc}")

    # R3: sha256 of the source's BYTES, exactly as they exist now -- not a
    # git commit SHA, and computed from the raw bytes (not a re-encoded
    # text form), so it matches "the source markdown's bytes" literally.
    sha_hex = hashlib.sha256(source_bytes).hexdigest()

    # I-B6/S13: a stateless read-and-compare against any prior render,
    # printed regardless of whether THIS invocation goes on to succeed.
    verdict = _fresh_stale_verdict(target_path, sha_hex)
    if verdict is not None:
        print(f"{deck}: {verdict}")

    # FR-015/R2: the ONE lazy-import site in this module. Everything pptx-
    # shaped below this point only runs once this succeeds.
    try:
        from pptx import Presentation
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        return Result(deck, OUTCOME_FAILED, None, "toolchain absent (python-pptx not installed)")

    kit = _PptxKit(Inches=Inches, Pt=Pt, PP_ALIGN=PP_ALIGN, MSO_AUTO_SIZE=MSO_AUTO_SIZE)

    try:
        # `Path.read_text()`'s universal-newline translation, replicated by
        # hand on the bytes already read above (deck_md.parse() documents
        # its input as "\n"-only line endings) -- avoids a second disk
        # read (and the TOCTOU it would invite) while matching
        # deck_md.parse_file()'s convention exactly.
        text = source_bytes.decode("utf-8").replace("\r\n", "\n").replace("\r", "\n")
        blocks = deck_md.parse(text)
        logical = _to_logical_slides(blocks)
        physical = _split_for_overflow(logical)
        h1_title = _plain_text(next(b.text for b in blocks if b.kind == deck_md.H1))
        prs = _assemble_presentation(
            Presentation, kit, h1_title, physical, source_path.as_posix(), sha_hex
        )
    except deck_md.DeckMdError as exc:
        return Result(deck, OUTCOME_FAILED, None, f"markdown outside supported constructs -- {exc}")
    except RenderError as exc:
        return Result(deck, OUTCOME_FAILED, None, str(exc))
    except UnicodeDecodeError as exc:
        return Result(deck, OUTCOME_FAILED, None, f"source markdown is not valid UTF-8: {exc}")
    except Exception as exc:  # noqa: BLE001 -- T10: anything else is a disclosed failure, never a crash
        return Result(deck, OUTCOME_FAILED, None, f"render transform failed: {exc}")

    # I-B3/O5: write to a temp file inside the target directory, then
    # os.replace() over the final path only on full success. The target is
    # never pre-deleted, so a mid-write failure leaves a prior good render
    # completely untouched.
    try:
        target_path.parent.mkdir(parents=True, exist_ok=True)
        fd, tmp_name = tempfile.mkstemp(
            prefix=f".{deck}.", suffix=".pptx.tmp", dir=str(target_path.parent)
        )
        os.close(fd)
    except OSError as exc:
        return Result(deck, OUTCOME_FAILED, None, f"cannot prepare the renders/ directory: {exc}")

    try:
        prs.save(tmp_name)
        os.replace(tmp_name, target_path)
    except Exception as exc:  # noqa: BLE001 -- mid-write failure: clean the temp file, leave target alone
        try:
            os.remove(tmp_name)
        except OSError:
            pass
        return Result(deck, OUTCOME_FAILED, None, f"write failed: {exc}")

    return Result(deck, OUTCOME_RENDERED, target_path, None)


# ---------------------------------------------------------------------------
# Disclosure + exit code (commands.md Sec 3 / Sec 4)
# ---------------------------------------------------------------------------


def _print_disclosure(feature_dir: Path, results: list[Result]) -> None:
    print(f"Deck render — {feature_dir.name}")
    label_width = max(len(r.deck) for r in results)
    for r in results:
        label = r.deck.ljust(label_width)
        if r.outcome == OUTCOME_RENDERED:
            print(f"  {label}   rendered   {r.path.as_posix()}")  # type: ignore[union-attr]
        elif r.outcome == OUTCOME_SKIPPED:
            print(f"  {label}   skipped    {r.reason}")
        else:
            print(f"  {label}   FAILED     {r.reason}")
    print()
    print("The markdown decks are unaffected and remain the artifact of record.")
    if any(r.outcome == OUTCOME_FAILED and "toolchain absent" in (r.reason or "") for r in results):
        print("Install the optional toolchain to render:  pip install python-pptx")


def _compute_exit_code(results: list[Result]) -> int:
    """The full six-outcome `both` matrix (commands.md Sec 4), stated as a
    general rule over any number of selected decks (1 or 2 in practice):
    all-skipped or no-failures -> 0; any failure alongside at least one
    rendered-or-skipped deck -> 2 (partial); every attempted deck failed
    (no rendered, no skipped) -> 4.
    """
    if not results:
        return EXIT_OK
    outcomes = [r.outcome for r in results]
    if all(o == OUTCOME_SKIPPED for o in outcomes):
        return EXIT_OK
    if any(o == OUTCOME_FAILED for o in outcomes):
        if any(o in (OUTCOME_RENDERED, OUTCOME_SKIPPED) for o in outcomes):
            return EXIT_PARTIAL
        return EXIT_ALL_FAILED
    return EXIT_OK


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


class _ArgParser(argparse.ArgumentParser):
    """argparse's default usage-error exit code is 2 -- which this
    contract already assigns the meaning "partial success" (commands.md
    Sec 4). Overridden so a bad invocation (an unrecognized flag, an
    out-of-choice deck argument) exits 3 ("invalid input") instead,
    matching this command's own convention rather than colliding with it.
    """

    def error(self, message):
        self.print_usage(sys.stderr)
        self.exit(EXIT_INVALID, f"{self.prog}: error: {message}\n")


def _build_arg_parser() -> _ArgParser:
    parser = _ArgParser(
        prog="render.py",
        description="Deterministic markdown -> pptx renderer for a feature's council "
        "defense decks (specs/006-deck-render/contracts/commands.md).",
    )
    parser.add_argument(
        "deck",
        nargs="?",
        choices=_CLI_DECK_CHOICES,
        default=None,
        metavar="{technical,overview,both}",
        help="Explicit deck selection -- overrides profile.yaml's deck_render entirely "
        "(FR-016), including when the profile says 'none'. Omit to render whatever the "
        "profile selects.",
    )
    parser.add_argument(
        "--feature",
        metavar="<dir>",
        default=None,
        help="Feature directory (contains council/defense-deck/, profile.yaml, renders/). "
        "Defaults to .specify/feature.json's feature_directory, else the current git branch.",
    )
    parser.add_argument(
        "--validate-profile",
        action="store_true",
        help="Validate the feature's deck_render profile key and exit; renders nothing. "
        "Exit 0 = valid (including absent). Exit 3 = out-of-enum or unreadable/unparseable YAML.",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    argv = sys.argv[1:] if argv is None else argv
    parser = _build_arg_parser()
    args = parser.parse_args(argv)

    try:
        feature_dir = _resolve_feature_dir(args.feature)
    except RenderInputError as exc:
        print(f"render.py: error: {exc}", file=sys.stderr)
        return EXIT_INVALID

    if args.validate_profile:
        try:
            resolve_deck_render(feature_dir=str(feature_dir))
        except ProfileKeyError as exc:
            print(f"deck_render: INVALID -- {exc}", file=sys.stderr)
            return EXIT_INVALID
        print(f"deck_render: valid ({feature_dir.as_posix()})")
        return EXIT_OK

    if args.deck:
        # FR-016/V4: an explicit deck argument wins outright -- the
        # profile is not consulted for selection at all when one is given.
        selected = _expand_selection(args.deck)
    else:
        try:
            profile_value = resolve_deck_render(feature_dir=str(feature_dir))
        except ProfileKeyError as exc:
            print(f"deck_render: INVALID -- {exc}", file=sys.stderr)
            return EXIT_INVALID
        selected = _expand_selection(profile_value)

    if not selected:
        # commands.md Sec 2 step 3: exactly one line, nothing else, exit 0.
        print(f"Deck render — {feature_dir.name}: nothing selected (deck_render: none). "
              "No file written.")
        return EXIT_OK

    ordered = tuple(d for d in RENDERABLE_DECKS if d in selected)
    results = [_render_deck(deck, feature_dir) for deck in ordered]
    _print_disclosure(feature_dir, results)
    return _compute_exit_code(results)


if __name__ == "__main__":
    sys.exit(main())
